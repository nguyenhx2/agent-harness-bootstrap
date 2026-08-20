#!/usr/bin/env python3
"""Execute one entry of route_sources.py's plan: read a source, write plain text.

Distilled from the docs-to-knowledge skill's scripts/ingest.py. Kept: the guarded optional
imports, the engine fallbacks, and the boundary markers that let a spec cite where a
statement came from. Dropped: that skill's bundle format (text.json, per-page assets, the
audit layer) - spec-builder needs readable text it can quote, not a knowledge bundle.

Output is one UTF-8 file at --out, with the boundaries the reader actually knows:

    === page 3 ===      pdf, via pymupdf
    === slide 7 ===     pptx, via python-pptx
    === sheet: Costs === xlsx, via openpyxl
    (docx has no page concept before rendering, so it carries none - saying "page 4" for a
     .docx would be a fabricated citation, which is the failure this whole skill guards.)

Every optional import is guarded and every failure names the fix. A source this cannot read
exits 2 with the reason: per reference/elicitation.md that becomes an OI-nn in section 11,
never a silent gap.

Usage:
    python ingest.py <source> --via native|anydoc|markitdown|read --out text.txt
    python ingest.py <scan.pdf> --via vision --render pages/
"""

from __future__ import annotations

import argparse
import os
import re
import shutil
import subprocess
import sys

TEXT_LAYER_MIN_CHARS_PER_PAGE = 80
IMAGE_EXT = (".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".tif", ".tiff")
PLAIN_EXT = (".md", ".markdown", ".txt", ".log")


class Unreadable(RuntimeError):
    """The source cannot be read with what is installed. Always carries the fix."""


def need(mod: str, pip: str, lost: str):
    """Import an optional reader or explain, in one line, what its absence costs."""
    try:
        return __import__(mod)
    except ImportError:
        raise Unreadable(f"{mod} is not installed. `pip install {pip}` {lost}")


# ------------------------------------------------------------------ engines
def via_anydoc(src: str) -> str:
    anydoc = need("anydoc", "firecrawl-anydoc",
                  "reads legacy Office, OpenDocument, RTF and spreadsheets as Markdown "
                  "tables offline; without it those route to markitdown or become an OI")
    try:
        return anydoc.to_markdown(src)
    except Exception as e:
        raise Unreadable(f"anydoc could not convert {os.path.basename(src)}: "
                         f"{type(e).__name__}: {e}. Try markitdown, or ask for a re-export.")


def via_markitdown(src: str) -> str:
    exe = shutil.which("markitdown")
    if exe:
        r = subprocess.run([exe, src], capture_output=True, text=True, encoding="utf-8")
        if r.returncode == 0 and (r.stdout or "").strip():
            return r.stdout
        err = (r.stderr or "").strip()
    else:
        err = "markitdown CLI not on PATH"
    try:
        from markitdown import MarkItDown
        return MarkItDown().convert(src).text_content
    except Exception as e:
        raise Unreadable(
            f"markitdown could not convert {os.path.basename(src)} ({err}; python API: "
            f"{e}). `pip install 'markitdown[all]'` converts HTML, CSV, JSON, XML, MSG, "
            f"EPUB, notebooks, archives and URLs; or invoke the `markitdown` skill.")


# ------------------------------------------------------------------ native readers
def read_pdf(src: str) -> str:
    fitz = need("fitz", "pymupdf",
                "reads PDF natively and keeps page boundaries; without it this PDF routes "
                "to markitdown (no page numbers) or to vision")
    doc = fitz.open(src)
    parts, chars = [], 0
    for i, page in enumerate(doc, 1):
        text = page.get_text().strip()
        chars += len(text)
        parts.append(f"=== page {i} ===\n{text}")
    n = max(doc.page_count, 1)
    doc.close()
    if chars / n < TEXT_LAYER_MIN_CHARS_PER_PAGE:
        raise Unreadable(
            f"{os.path.basename(src)} yielded {chars} characters over {n} page(s) "
            f"({chars / n:.1f}/page, under {TEXT_LAYER_MIN_CHARS_PER_PAGE}): this is a "
            f"scan, and the text layer is a lie. Re-run with `--via vision --render <dir>` "
            f"and read the pages, or ask the user for a text-bearing export. Do NOT write "
            f"requirements from what little came out.")
    return "\n\n".join(parts)


def read_docx(src: str) -> str:
    need("docx", "python-docx",
         "reads .docx paragraphs and tables; without it this file routes to anydoc or "
         "markitdown, which flatten table structure")
    import docx
    d = docx.Document(src)
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for t, table in enumerate(d.tables, 1):
        parts.append(f"=== table {t} ===")
        for row in table.rows:
            vals = [c.text.strip() for c in row.cells]
            if any(vals):
                parts.append("\t".join(vals))
    return "\n".join(parts)


def read_xlsx(src: str) -> str:
    need("openpyxl", "openpyxl",
         "reads .xlsx sheets; without it this workbook routes to anydoc or markitdown, "
         "which give Markdown tables instead of cell rows")
    import openpyxl
    wb = openpyxl.load_workbook(src, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"=== sheet: {ws.title} ===")
        for row in ws.iter_rows(values_only=True):
            vals = ["" if v is None else str(v) for v in row]
            if any(v.strip() for v in vals):
                parts.append("\t".join(vals).rstrip("\t"))
    wb.close()
    return "\n".join(parts)


def read_pptx(src: str) -> str:
    need("pptx", "python-pptx",
         "reads .pptx slides, shapes and table cells; without it this deck routes to "
         "anydoc or markitdown, which lose slide boundaries")
    from pptx import Presentation

    def walk(shapes, out):
        for sh in shapes:
            if sh.shape_type == 6:  # group
                walk(sh.shapes, out)
            elif getattr(sh, "has_table", False):
                for row in sh.table.rows:
                    vals = [c.text_frame.text.strip() for c in row.cells]
                    if any(vals):
                        out.append("\t".join(vals))
            elif getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip():
                out.append(sh.text_frame.text.strip())

    parts = []
    for i, slide in enumerate(Presentation(src).slides, 1):
        lines: list[str] = []
        walk(slide.shapes, lines)
        parts.append(f"=== slide {i} ===\n" + "\n".join(lines))
    return "\n\n".join(parts)


def read_plain(src: str) -> str:
    with open(src, encoding="utf-8", errors="replace") as f:
        return f.read()


NATIVE = {".pdf": read_pdf, ".docx": read_docx, ".xlsx": read_xlsx, ".xlsm": read_xlsx,
          ".pptx": read_pptx}


# ------------------------------------------------------------------ vision
def render_pages(src: str, out_dir: str) -> list[str]:
    """Rasterize, ONLY for the vision route. A readable document is never rendered here:
    docs-to-knowledge renders everything because it cross-checks diagrams against text;
    a BA quoting a readable PDF gains nothing from a PNG of it."""
    os.makedirs(out_dir, exist_ok=True)
    if src.lower().endswith(IMAGE_EXT):
        dest = os.path.join(out_dir, os.path.basename(src))
        shutil.copy(src, dest)
        return [dest]
    fitz = need("fitz", "pymupdf",
                "rasterizes PDF pages for a vision read; without it a scanned PDF cannot "
                "be read at all and must become an OI-nn")
    doc = fitz.open(src)
    made = []
    for i, page in enumerate(doc, 1):
        dest = os.path.join(out_dir, f"p{i:03d}.png")
        page.get_pixmap(matrix=fitz.Matrix(2, 2)).save(dest)
        made.append(dest)
    doc.close()
    return made


# ------------------------------------------------------------------ driver
def extract(src: str, via: str) -> str:
    if re.match(r"^https?://", src, re.I):
        return via_markitdown(src)
    ext = os.path.splitext(src)[1].lower()
    if via == "read" or (via == "auto" and ext in PLAIN_EXT):
        return read_plain(src)
    if via == "anydoc":
        return via_anydoc(src)
    if via == "markitdown":
        return via_markitdown(src)
    reader = NATIVE.get(ext)
    if reader is None:
        raise Unreadable(
            f"no native reader for `{ext or 'a file with no extension'}` in this skill. "
            f"route_sources.py routes it to anydoc or markitdown; if neither is installed, "
            f"ask the user to re-export as PDF, DOCX, XLSX, PPTX or Markdown.")
    return reader(src)


def main() -> int:
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
        sys.stderr.reconfigure(encoding="utf-8", errors="replace")
    ap = argparse.ArgumentParser(description="Read one source into plain text.")
    ap.add_argument("source", help="file path or http(s) URL")
    ap.add_argument("--via", default="auto",
                    choices=("auto", "native", "read", "anydoc", "markitdown", "vision"),
                    help="who owns the text layer (route_sources.py decides this)")
    ap.add_argument("--out", help="write the extracted text here")
    ap.add_argument("--render", help="vision route only: write page images here")
    a = ap.parse_args()

    try:
        if a.via == "vision":
            if not a.render:
                sys.stderr.write("--via vision needs --render <dir>: there is no text to "
                                 "write, only pages for a vision agent to read.\n")
                return 2
            made = render_pages(a.source, a.render)
            print(f"{os.path.basename(a.source)}: no trustworthy text layer. "
                  f"{len(made)} page image(s) -> {a.render}")
            print("  Read them with a vision agent and say in the spec that the content "
                  "was read from pixels.")
            return 0
        text = extract(a.source, a.via)
    except Unreadable as e:
        sys.stderr.write(f"UNREADABLE {os.path.basename(a.source) or a.source}: {e}\n")
        sys.stderr.write("  This is an OI-nn in section 11 naming the file, why it cannot "
                         "be read, and the fix above. Never a silent gap.\n")
        return 2

    if not a.out:
        sys.stdout.write(text)
        return 0
    parent = os.path.dirname(os.path.abspath(a.out))
    if parent:
        os.makedirs(parent, exist_ok=True)
    with open(a.out, "w", encoding="utf-8", newline="\n") as f:
        f.write(text)
    stripped = text.strip()
    print(f"{os.path.basename(a.source) or a.source}: {len(stripped)} chars via {a.via} "
          f"-> {a.out}")
    if not stripped:
        print("  WARNING: the reader returned nothing. An empty extraction is not an "
              "empty document - treat it as unreadable and raise an OI-nn.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
